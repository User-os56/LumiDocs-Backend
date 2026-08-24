import os

import pypdf
import docx
from pptx import Presentation


def extract_text_from_file(file_path):
    """
    Extract plain text from supported document formats.

    Supported:
    - PDF
    - DOCX
    - PPTX

    Returns:
        str: Extracted text.
    """

    ext = os.path.splitext(file_path)[1].lower()

    try:

        # =====================================================
        # PDF
        # =====================================================

        if ext == ".pdf":

            reader = pypdf.PdfReader(file_path)

            pages = []

            for page_number, page in enumerate(
                reader.pages,
                start=1
            ):

                extracted = page.extract_text()

                if extracted:
                    pages.append(
                        f"\n[Page {page_number}]\n"
                        f"{extracted.strip()}"
                    )

            return "\n".join(pages).strip()


        # =====================================================
        # DOCX
        # =====================================================

        elif ext == ".docx":

            document = docx.Document(file_path)

            paragraphs = []

            for paragraph in document.paragraphs:

                text = paragraph.text.strip()

                if text:
                    paragraphs.append(text)

            # Also extract text from tables
            for table in document.tables:

                for row in table.rows:

                    row_text = []

                    for cell in row.cells:

                        cell_text = cell.text.strip()

                        if cell_text:
                            row_text.append(cell_text)

                    if row_text:
                        paragraphs.append(
                            " | ".join(row_text)
                        )

            return "\n".join(paragraphs).strip()


        # =====================================================
        # PPTX
        # =====================================================

        elif ext == ".pptx":

            presentation = Presentation(file_path)

            slides_text = []

            for slide_number, slide in enumerate(
                presentation.slides,
                start=1
            ):

                slide_content = []

                for shape in slide.shapes:

                    if not hasattr(shape, "text"):
                        continue

                    text = shape.text.strip()

                    if text:
                        slide_content.append(text)

                if slide_content:

                    slides_text.append(
                        f"\n[Slide {slide_number}]\n"
                        + "\n".join(slide_content)
                    )

            return "\n".join(slides_text).strip()


        # =====================================================
        # Unsupported format
        # =====================================================

        else:

            raise ValueError(
                f"Unsupported file format: {ext}"
            )

    except Exception as e:

        print(
            f"Error extracting text from "
            f"{file_path}: {type(e).__name__}: {e}"
        )

        raise